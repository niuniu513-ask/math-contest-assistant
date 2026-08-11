#!/usr/bin/env python3
"""物理示意图模板：用 python-pptx 以代码生成可复现的物理示意图。

适用：坐标/光线/矢量/角度/受力/场/几何布局等示意性插图（如太阳反射、
定日镜场、受力分析、坐标系）。这类图是"几何关系示意"，用 PPT 自由绘制
比 Visio 工程图更合适，且可纳入复现清单。

用法：
    from schematic_pptx import Schematic
    s = Schematic()
    s.axes((0, 0), (8, 0), (0, 5), xlabel="x (m)", ylabel="y (m)")
    s.arrow((1, 1), (4, 3), color="E69F00", width=2.5, label="太阳光线")
    s.ray((4, 1), 120, length=3)          # 从点 (4,1) 沿 120° 方向
    s.arc((4, 3), 0.8, 30, 75, label="θ")
    s.text((4, 4.2), "集热器中心", fontsize=10)
    s.mirror((4, 1), width=2, angle=-20)  # 画一个镜面线段
    s.save("示意图.pptx", export_png="示意图.png")

依赖：pip install python-pptx
配色：建议用 plot_style.PALETTE（#0072B2 等），见 物理示意图绘制规范.md
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

# ---- 画布：逻辑坐标 -> 英寸/EMU ----
EMU_PER_INCH = 914400
SLIDE_W = 13.333        # 16:9 英寸
SLIDE_H = 7.5


class Schematic:
    def __init__(self, xlim=(0, 12), ylim=(0, 6), slide_w=SLIDE_W, slide_h=SLIDE_H):
        self.prs = Presentation()
        self.prs.slide_width = Inches(slide_w)
        self.prs.slide_height = Inches(slide_h)
        self.slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # 空白版式
        self.xlim = xlim
        self.ylim = ylim
        # 页面留边（英寸），使绘图区居中
        self.margin_x = 0.6
        self.margin_y = 0.6
        self.plot_w_in = slide_w - 2 * self.margin_x
        self.plot_h_in = slide_h - 2 * self.margin_y

    def _to_emu(self, x, y):
        """逻辑坐标 -> (left_emu, top_emu) 的左上角（python-pptx 用左上角定位）。"""
        # 逻辑 x -> 英寸：映射 xlim 到绘图区
        xi = self.margin_x + (x - self.xlim[0]) / (self.xlim[1] - self.xlim[0]) * self.plot_w_in
        # 逻辑 y（向上）-> 英寸 y（向下）
        yi = self.margin_y + (self.ylim[1] - y) / (self.ylim[1] - self.ylim[0]) * self.plot_h_in
        return Emu(int(xi * EMU_PER_INCH)), Emu(int(yi * EMU_PER_INCH))

    def _span_emu(self, dx, dy):
        """逻辑长度 -> (width_emu, height_emu)。"""
        w = dx / (self.xlim[1] - self.xlim[0]) * self.plot_w_in
        h = dy / (self.ylim[1] - self.ylim[0]) * self.plot_h_in
        return Emu(int(abs(w) * EMU_PER_INCH)), Emu(int(abs(h) * EMU_PER_INCH))

    @staticmethod
    def _rgb(color: str) -> RGBColor:
        color = color.lstrip("#")
        return RGBColor(int(color[0:2], 16), int(color[2:4], 16), int(color[4:6], 16))

    # ---- 原语 ----
    def line(self, p1, p2, color="222222", width=1.5, dash=None):
        """直线段。p1,p2=(x,y)。dash='dash'/'dot' 可选。"""
        x1, y1 = self._to_emu(*p1)
        x2, y2 = self._to_emu(*p2)
        conn = self.slide.shapes.add_connector(1, x1, y1, x2, y2)  # 1=straight
        conn.line.color.rgb = self._rgb(color)
        conn.line.width = Pt(width)
        if dash:
            ln = conn.line._get_or_add_ln()
            prst = ln.find(qn("a:prstDash"))
            if prst is None:
                prst = ln.makeelement(qn("a:prstDash"), {})
                ln.insert(0, prst)
            prst.set("val", dash)
        return conn

    def arrow(self, p1, p2, color="222222", width=2.0, label=None, label_pos=0.5,
              label_offset=(0.15, 0.12), fontsize=10):
        """带箭头线段 + 可选标签。箭头在 p2 端。"""
        x1, y1 = self._to_emu(*p1)
        x2, y2 = self._to_emu(*p2)
        conn = self.slide.shapes.add_connector(1, x1, y1, x2, y2)
        conn.line.color.rgb = self._rgb(color)
        conn.line.width = Pt(width)
        conn.line.end_arrowhead = True
        if label:
            lx = p1[0] + (p2[0] - p1[0]) * label_pos + label_offset[0]
            ly = p1[1] + (p2[1] - p1[1]) * label_pos + label_offset[1]
            self.text((lx, ly), label, fontsize=fontsize, color=color)
        return conn

    def ray(self, p, angle_deg, length, color="222222", width=2.0, label=None):
        """从点 p 沿 angle_deg（度，x 正向为 0，逆时针）画长度 length 的箭头射线。"""
        import math
        a = math.radians(angle_deg)
        p2 = (p[0] + length * math.cos(a), p[1] + length * math.sin(a))
        return self.arrow(p, p2, color=color, width=width, label=label)

    def axis(self, start, end, color="444444", width=1.5, label=None, label_offset=(0.15, -0.15)):
        """坐标轴：带箭头，可带轴名标签。"""
        return self.arrow(start, end, color=color, width=width, label=label, label_offset=label_offset)

    def axes(self, origin, x_end, y_end, xlabel="x", ylabel="y", color="444444", width=1.5):
        """画 x/y 坐标轴（x 右、y 上）。"""
        self.axis(origin, x_end, color=color, width=width, label=xlabel, label_offset=(0.1, -0.18))
        self.axis(origin, y_end, color=color, width=width, label=ylabel, label_offset=(-0.35, 0.05))

    def arc(self, center, radius, start_deg, end_deg, color="D55E00", width=1.5,
            label=None, label_fontsize=10):
        """角度弧（圆弧线），用于标注角度。start/end_deg 度。"""
        import math
        n = 32
        pts = []
        for i in range(n + 1):
            a = math.radians(start_deg + (end_deg - start_deg) * i / n)
            pts.append((center[0] + radius * math.cos(a), center[1] + radius * math.sin(a)))
        for i in range(n):
            self.line(pts[i], pts[i + 1], color=color, width=width)
        if label:
            a = math.radians((start_deg + end_deg) / 2)
            lp = (center[0] + (radius + 0.18) * math.cos(a), center[1] + (radius + 0.18) * math.sin(a))
            self.text(lp, label, fontsize=label_fontsize, color=color)
        return pts

    def mirror(self, center, width, angle_deg, color="0072B2", thick=0.10, label=None):
        """画一个旋转的镜面矩形（宽度 width，表面方向角 angle_deg，厚度 thick）。

        用自由多边形按真实旋转角绘制，避免 add_shape+rotation 的双旋转问题。
        """
        import math
        a = math.radians(angle_deg)
        hw = width / 2.0
        hh = thick / 2.0
        sx, sy = math.cos(a), math.sin(a)                       # 镜面方向
        tx, ty = math.cos(a + math.pi / 2), math.sin(a + math.pi / 2)  # 厚度方向
        cx, cy = center
        corners = [
            (cx + hw * sx - hh * tx, cy + hw * sy - hh * ty),
            (cx - hw * sx - hh * tx, cy - hw * sy - hh * ty),
            (cx - hw * sx + hh * tx, cy - hw * sy + hh * ty),
            (cx + hw * sx + hh * tx, cy + hw * sy + hh * ty),
        ]
        pts = [self._to_emu(*c) for c in corners]
        fb = self.slide.shapes.build_freeform(pts[0][0], pts[0][1])
        fb.add_line_segments([(p[0], p[1]) for p in pts[1:]], close=True)
        shape = fb.convert_to_shape()
        shape.fill.solid(); shape.fill.fore_color.rgb = self._rgb(color)
        shape.line.color.rgb = self._rgb(color)
        if label:
            self.text((center[0], center[1] - thick), label, fontsize=10, color=color)
        return shape

    def tower(self, base, height, width=0.4, color="37474F", label=None):
        """画一根竖直柱体（塔/杆）。base=(x, y底)。"""
        l, t = self._to_emu(base[0] - width / 2, base[1] + height)
        w, h = self._span_emu(width, height)
        shape = self.slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
        shape.fill.solid(); shape.fill.fore_color.rgb = self._rgb(color)
        shape.line.color.rgb = self._rgb(color)
        if label:
            self.text((base[0] + width, base[1] + height * 0.8), label, fontsize=10, color=color)
        return shape

    def text(self, pos, text, fontsize=10, color="222222", bold=False, align=PP_ALIGN.LEFT):
        """文本框标注。pos=(x,y) 为文本锚点。"""
        l, t = self._to_emu(*pos)
        box = self.slide.shapes.add_textbox(l, t, Inches(2.2), Inches(0.4))
        tf = box.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run(); run.text = text
        run.font.size = Pt(fontsize)
        run.font.bold = bold
        run.font.color.rgb = self._rgb(color)
        return box

    def circle(self, center, radius, color="222222", fill=None, width=1.5):
        """圆点/圆环。fill='#' 填充色可选。"""
        l, t = self._to_emu(center[0] - radius, center[1] + radius)
        w, h = self._span_emu(2 * radius, 2 * radius)
        shape = self.slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, w, h)
        if fill:
            shape.fill.solid(); shape.fill.fore_color.rgb = self._rgb(fill)
        else:
            shape.fill.background()
        shape.line.color.rgb = self._rgb(color)
        shape.line.width = Pt(width)
        return shape

    def save(self, path, export_png=None, export_svg=None):
        """保存 pptx；可选经 Microsoft PowerPoint COM 导出 PNG/SVG。

        参数 export_png/export_svg：True 用与 path 同名默认扩展名；传字符串则
        指定路径；None 不导出该格式。导出依赖本机 PowerPoint（无 PPT 环境
        只存 pptx 并标注 BLOCKED）。
        """
        out = Path(path)
        self.prs.save(str(out))
        do_png = export_png is True or isinstance(export_png, str)
        do_svg = export_svg is True or isinstance(export_svg, str)
        if do_png or do_svg:
            png = str(out.with_suffix(".png")) if export_png is True else (str(export_png) if isinstance(export_png, str) else None)
            svg = str(out.with_suffix(".svg")) if export_svg is True else (str(export_svg) if isinstance(export_svg, str) else None)
            self._export_com(out, png, svg)
        return out

    def _export_com(self, pptx_path, png_path, svg_path):
        """用 PowerPoint COM 导出 PNG（约 300dpi）。PPTX 本身即矢量源，
        需要 EMF/SVG 时在 PowerPoint 中手动另存（自动导出会生成资源文件夹）。"""
        try:
            import win32com.client  # pywin32
            import os
            app = win32com.client.Dispatch("PowerPoint.Application")
            prs = app.Presentations.Open(str(os.path.abspath(pptx_path)), WithWindow=False)
            if png_path and png_path != "None":
                for i in range(1, prs.Slides.Count + 1):
                    prs.Slides(i).Export(str(os.path.abspath(png_path)), "PNG", 3300, 1856)
            prs.Close()
            app.Quit()
            print(f"[schematic_pptx] 已导出 PNG: {png_path}")
        except Exception as exc:
            print(f"[schematic_pptx] COM 导出失败（仅保留 pptx）: {exc}")


# ---- 示例：定日镜反射示意（2023 CUMCM A）----
def demo_heliostat(path="示意图_定日镜反射.pptx", export_png=True, export_svg=False):
    import math
    s = Schematic(xlim=(0, 12), ylim=(0, 6))
    M = (5.2, 1.0)      # 镜面中心
    S = (2.0, 5.2)      # 太阳
    R = (9.5, 3.5)      # 集热器中心
    # 反射几何：法向 n = normalize(unit(M→S) + unit(M→R))
    def unit(a, b):
        dx, dy = b[0] - a[0], b[1] - a[1]
        L = math.hypot(dx, dy)
        return (dx / L, dy / L)
    uMS = unit(M, S); uMR = unit(M, R)
    nx, ny = uMS[0] + uMR[0], uMS[1] + uMR[1]
    L = math.hypot(nx, ny); nx, ny = nx / L, ny / L
    n_deg = math.degrees(math.atan2(ny, nx))      # 法向仰角
    mirror_tilt = n_deg - 90.0                     # 镜面垂直于法向
    theta = abs(math.degrees(math.atan2(ny, nx)) - math.degrees(math.atan2(uMS[1], uMS[0])))
    alpha_s = math.degrees(math.atan2(S[1] - M[1], S[0] - M[0]))  # 太阳方向角（向上）
    sun_ray_deg = math.degrees(math.atan2(uMS[1], uMS[0]))        # 入射光方向角

    # 地面
    s.line((0, 0.3), (12, 0.3), color="8D99AE", width=1.5)
    # 坐标系
    s.axes((0.8, 0.5), (3.2, 0.5), (0.8, 2.6), xlabel="x (m)", ylabel="y (m)")
    # 塔与集热器
    s.tower((9.5, 0.3), 3.2, width=0.5, label="吸收塔")
    s.circle(R, 0.35, color="D55E00", fill="D55E00")
    s.text((9.3, 4.15), "集热器中心", fontsize=10)
    # 太阳
    s.circle(S, 0.3, color="E69F00", fill="E69F00")
    s.text((2.3, 5.45), "太阳", fontsize=10, color="E69F00")
    # 定日镜（镜面垂直于法向 n）
    s.mirror(M, 2.4, angle_deg=mirror_tilt, label="定日镜")
    # 入射光（太阳→镜心）
    s.arrow((2.3, 4.9), M, color="E69F00", width=2.5, label="入射光", label_offset=(0.25, 0.25))
    # 反射光（镜心→集热器）
    s.arrow(M, R, color="009E73", width=2.5, label="反射光", label_offset=(0.25, 0.3))
    # 法向虚线（从镜心沿 n 延伸，灰色以与光线区分）
    n_end = (M[0] + 2.6 * nx, M[1] + 2.6 * ny)
    s.line(M, n_end, color="555555", width=1.2, dash="dash")
    s.text((n_end[0] + 0.15, n_end[1]), "法向 n", fontsize=9, color="555555")
    # 入射角 θ（入射光与法向之间）
    s.arc(M, 0.9, min(sun_ray_deg, n_deg), max(sun_ray_deg, n_deg), color="6A5ACD", label="θ")
    # 太阳高度角 αs（水平线与入射光之间，画在镜心附近）
    s.arc(M, 0.55, min(sun_ray_deg, 180.0), max(sun_ray_deg, 180.0), color="E69F00", label="αs")
    s.text((M[0] + 0.2, M[1] + 0.5), f"αs≈{180-sun_ray_deg:.0f}°, θ≈{theta:.0f}°", fontsize=9)
    s.text((0.6, 0.05), "图：塔式光热定日镜反射几何示意", fontsize=11, bold=True)
    s.save(path, export_png=export_png, export_svg=export_svg)
    print(f"[schematic_pptx] 已保存: {path}  (αs={180-sun_ray_deg:.1f}°, θ={theta:.1f}°, 镜面倾角={mirror_tilt:.1f}°)")


if __name__ == "__main__":
    import sys
    if len(sys.argv) > 1 and sys.argv[1] in ("-h", "--help"):
        print("用法: python schematic_pptx.py [输出路径]")
        print("不传参数时运行内置示例并生成 示意图_定日镜反射.pptx/.png。")
        sys.exit(0)
    demo_heliostat(sys.argv[1] if len(sys.argv) > 1 else "示意图_定日镜反射.pptx")
